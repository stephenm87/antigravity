#!/usr/bin/env python3
"""
build-day8-pptx.py
Builds the Day 8 Chinese Civil War presentation as a PowerPoint file (.pptx).
Uses python-pptx with a dark theme aesthetic matching the HTML version.

Usage: python3 build-day8-pptx.py

Requires: pip3 install python-pptx requests
"""

import os
import io
import tempfile
import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "day8-ccw-presentation.pptx")

# ─── Color Palette ──────────────────────────────────────────
BG_DARK      = RGBColor(0x1A, 0x1A, 0x1E)
BG_SLIDE     = RGBColor(0x22, 0x22, 0x28)
TEXT_PRIMARY  = RGBColor(0xF0, 0xEB, 0xE3)
TEXT_SECONDARY= RGBColor(0xC0, 0xBB, 0xB2)
ACCENT_RED    = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_RED_LT = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_GOLD   = RGBColor(0xD4, 0xA8, 0x53)
ACCENT_GOLD_LT= RGBColor(0xF0, 0xD6, 0x8A)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x00, 0x00, 0x00)

# Connections colors
CONN_YELLOW   = RGBColor(0xF1, 0xC4, 0x0F)
CONN_GREEN    = RGBColor(0x7D, 0xCE, 0xA0)
CONN_BLUE     = RGBColor(0x6F, 0xA8, 0xDC)
CONN_PURPLE   = RGBColor(0xBB, 0x8F, 0xCE)
CONN_UNSOLVED = RGBColor(0x3A, 0x3A, 0x42)

# Work Time card colors
WT_ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
WT_GRAY       = RGBColor(0x7F, 0x8C, 0x8D)
WT_GOLD       = RGBColor(0xF1, 0xC4, 0x0F)
WT_BLUE       = RGBColor(0x5D, 0xAD, 0xE2)
WT_GREEN      = RGBColor(0x27, 0xAE, 0x60)

# ─── Image URLs (Wikimedia Commons, Public Domain) ──────────
IMG_MAO_PROCLAIM = "https://upload.wikimedia.org/wikipedia/commons/8/85/Mao_Proclaiming_New_China.JPG"
IMG_CCW_MAP = "https://upload.wikimedia.org/wikipedia/commons/2/2e/Chinese_civil_war_map_02.jpg"

# ─── Helpers ────────────────────────────────────────────────
def download_image(url):
    """Download image and return as BytesIO."""
    print(f"  Downloading: {url.split('/')[-1]}...")
    headers = {
        'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
    }
    resp = requests.get(url, timeout=30, headers=headers)
    resp.raise_for_status()
    return io.BytesIO(resp.content)

def set_slide_bg(slide, color):
    """Set slide background to a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bar(slide, y_pos, width, color=ACCENT_RED, height=Pt(4)):
    """Add a thin colored accent bar to the slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y_pos, width, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_gold_underline(slide, x, y, width=Inches(1.2)):
    """Add a gold underline accent."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, width, Pt(3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GOLD
    bar.line.fill.background()
    return bar

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=TEXT_PRIMARY, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(text_frame, items, font_size=16, font_color=TEXT_SECONDARY,
                    level=0, font_name='Calibri', bold=False):
    """Add bulleted items to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.font.bold = bold
        p.level = level
        p.space_before = Pt(4)
        p.space_after = Pt(4)

def add_card(slide, left, top, width, height, bg_color, border_color=None):
    """Add a rounded rectangle card."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

def add_conn_cell(slide, col, row, text, color, text_color=BLACK,
                  grid_left=Inches(3.3), grid_top=Inches(1.8),
                  cell_w=Inches(1.85), cell_h=Inches(1.0), gap=Inches(0.12)):
    """Add a connections puzzle cell."""
    x = grid_left + col * (cell_w + gap)
    y = grid_top + row * (cell_h + gap)
    
    cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cell_w, cell_h)
    cell.fill.solid()
    cell.fill.fore_color.rgb = color
    cell.line.fill.background()
    
    tf = cell.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = 'Calibri'
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ─── Slide Building ─────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
blank_layout = prs.slide_layouts[6]  # blank layout

# ═══════════════ SLIDE 1: TITLE ═══════════════
print("Building Slide 1: Title...")
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, BG_DARK)

# Top accent gradient bar (simplified to single color) — added last for z-order
# We'll add it after setting up the background

# Background image (semi-transparent overlay achieved by dark overlay shape)
try:
    img_data = download_image(IMG_MAO_PROCLAIM)
    slide1.shapes.add_picture(img_data, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Dark overlay with transparency — set to 40% so image shows through
    overlay = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BG_DARK
    overlay.line.fill.background()
    # Set transparency via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    fill_elem = overlay._element.find('.//' + qn('a:solidFill'))
    if fill_elem is not None:
        srgbClr = fill_elem.find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '40000')  # 40% opacity — lets image show through clearly
except Exception as e:
    print(f"  Warning: Could not add background image: {e}")

# Top accent bar (on top of overlay)
add_accent_bar(slide1, Inches(0), SLIDE_W, ACCENT_RED)

# "DAY 8" badge
badge = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.0), Inches(2.3), Inches(0.5)
)
badge.fill.solid()
badge.fill.fore_color.rgb = ACCENT_RED
badge.line.fill.background()
tf = badge.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = "DAY 8"
p.font.size = Pt(16)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = 'Calibri'
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Title text
add_text_box(slide1, Inches(1.5), Inches(2.8), Inches(10.3), Inches(2.0),
             "Fall of the Nationalists and\nFounding the People's Republic of China",
             font_size=38, font_color=TEXT_PRIMARY, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Georgia')

# Date range
add_text_box(slide1, Inches(4), Inches(4.8), Inches(5.3), Inches(0.8),
             "1946 – 49", font_size=32, font_color=ACCENT_GOLD, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Georgia')

# Subtitle
add_text_box(slide1, Inches(4), Inches(5.7), Inches(5.3), Inches(0.5),
             "March 17/A  –  March 18/B", font_size=14, font_color=TEXT_SECONDARY,
             alignment=PP_ALIGN.CENTER)

# ═══════════════ SLIDE 2: AGENDA ═══════════════
print("Building Slide 2: Agenda...")
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, BG_SLIDE)
add_accent_bar(slide2, Inches(0), SLIDE_W, ACCENT_RED)

# Title
add_text_box(slide2, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
             "Day 8 – Fall of Nationalists and Founding of the People's Republic of China",
             font_size=28, font_color=TEXT_PRIMARY, bold=True, font_name='Georgia')
add_gold_underline(slide2, Inches(0.7), Inches(1.1))

# Objective card
obj_card = add_card(slide2, Inches(0.7), Inches(1.6), Inches(5.8), Inches(2.2),
                    RGBColor(0x28, 0x28, 0x30), RGBColor(0x3A, 0x3A, 0x42))
# Red left border for Objective
add_card(slide2, Inches(0.7), Inches(1.6), Pt(5), Inches(2.2), ACCENT_RED)

obj_title = add_text_box(slide2, Inches(1.0), Inches(1.75), Inches(5.2), Inches(0.4),
                         "OBJECTIVE", font_size=16, font_color=ACCENT_GOLD, bold=True)
obj_body = slide2.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(1.4))
tf = obj_body.text_frame
tf.word_wrap = True
add_bullet_list(tf, [
    "I can analyze multiple and complex causes of the success of the CCP in the Chinese Civil War.",
    "I can develop a precise and knowledgeable thesis statement"
], font_size=14, font_color=TEXT_SECONDARY)

# Agenda card
agenda_card = add_card(slide2, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.2),
                       RGBColor(0x28, 0x28, 0x30), RGBColor(0x3A, 0x3A, 0x42))
add_card(slide2, Inches(6.8), Inches(1.6), Pt(5), Inches(2.2), ACCENT_GOLD)

add_text_box(slide2, Inches(7.1), Inches(1.75), Inches(5.2), Inches(0.4),
             "AGENDA", font_size=16, font_color=ACCENT_GOLD, bold=True)
agenda_body = slide2.shapes.add_textbox(Inches(7.1), Inches(2.2), Inches(5.2), Inches(1.4))
tf = agenda_body.text_frame
tf.word_wrap = True
add_bullet_list(tf, [
    "Warm-up: CCW Connections",
    "Notes: The impact of WWII on the CCW",
    "Developing a Complex Thesis statement Review",
    "What Caused the Success of the CCP Document Set"
], font_size=14, font_color=TEXT_SECONDARY)

# Homework card (full width)
hw_card = add_card(slide2, Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.5),
                   RGBColor(0x28, 0x28, 0x30), RGBColor(0x3A, 0x3A, 0x42))
add_card(slide2, Inches(0.7), Inches(4.2), Pt(5), Inches(2.5), ACCENT_GOLD_LT)

add_text_box(slide2, Inches(1.0), Inches(4.35), Inches(5.2), Inches(0.4),
             "HOMEWORK", font_size=16, font_color=ACCENT_GOLD, bold=True)
hw_body = slide2.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.5))
tf = hw_body.text_frame
tf.word_wrap = True
add_bullet_list(tf, [
    "Have your source for Entry 5 ready for next class (Reassessment)",
    "Next two classes will workdays",
    "Final Project due at the start of class on April 9/10"
], font_size=14, font_color=TEXT_SECONDARY)


# ═══════════════ SLIDES 3–7: CONNECTIONS ═══════════════
# Define the grid states
connections_states = [
    # Slide 3: All unsolved
    {
        "rows": [
            [("Gangsters", CONN_UNSOLVED, TEXT_PRIMARY), ("Foreigners", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Comintern", CONN_UNSOLVED, TEXT_PRIMARY), ("Kuomintang", CONN_UNSOLVED, TEXT_PRIMARY)],
            [("Capitalists", CONN_UNSOLVED, TEXT_PRIMARY), ("GMD", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Peasants", CONN_UNSOLVED, TEXT_PRIMARY), ("Land-owners", CONN_UNSOLVED, TEXT_PRIMARY)],
            [("Gongchandang", CONN_UNSOLVED, TEXT_PRIMARY), ("Workers", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Communists", CONN_UNSOLVED, TEXT_PRIMARY), ("Guomindang", CONN_UNSOLVED, TEXT_PRIMARY)],
            [("Nationalists", CONN_UNSOLVED, TEXT_PRIMARY), ("The Chinese\nCommunist Party", CONN_UNSOLVED, TEXT_PRIMARY),
             ("CCP", CONN_UNSOLVED, TEXT_PRIMARY), ("Soviet Union", CONN_UNSOLVED, TEXT_PRIMARY)]
        ]
    },
    # Slide 4: Yellow solved
    {
        "rows": [
            [("CCP", CONN_YELLOW, BLACK), ("Communists", CONN_YELLOW, BLACK),
             ("Gongchandang", CONN_YELLOW, BLACK), ("The Chinese\nCommunist Party", CONN_YELLOW, BLACK)],
            [("Capitalists", CONN_UNSOLVED, TEXT_PRIMARY), ("GMD", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Peasants", CONN_UNSOLVED, TEXT_PRIMARY), ("Land-owners", CONN_UNSOLVED, TEXT_PRIMARY)],
            [("Comintern", CONN_UNSOLVED, TEXT_PRIMARY), ("Workers", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Foreigners", CONN_UNSOLVED, TEXT_PRIMARY), ("Guomindang", CONN_UNSOLVED, TEXT_PRIMARY)],
            [("Nationalists", CONN_UNSOLVED, TEXT_PRIMARY), ("Kuomintang", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Gangsters", CONN_UNSOLVED, TEXT_PRIMARY), ("Soviet Union", CONN_UNSOLVED, TEXT_PRIMARY)]
        ]
    },
    # Slide 5: Yellow + Green
    {
        "rows": [
            [("CCP", CONN_YELLOW, BLACK), ("Communists", CONN_YELLOW, BLACK),
             ("Gongchandang", CONN_YELLOW, BLACK), ("The Chinese\nCommunist Party", CONN_YELLOW, BLACK)],
            [("Nationalists", CONN_GREEN, BLACK), ("GMD", CONN_GREEN, BLACK),
             ("Kuomintang", CONN_GREEN, BLACK), ("Guomindang", CONN_GREEN, BLACK)],
            [("Peasants", CONN_UNSOLVED, TEXT_PRIMARY), ("Foreigners", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Comintern", CONN_UNSOLVED, TEXT_PRIMARY), ("Capitalists", CONN_UNSOLVED, TEXT_PRIMARY)],
            [("Gangsters", CONN_UNSOLVED, TEXT_PRIMARY), ("Workers", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Land-owners", CONN_UNSOLVED, TEXT_PRIMARY), ("Soviet Union", CONN_UNSOLVED, TEXT_PRIMARY)]
        ]
    },
    # Slide 6: Yellow + Green + Blue
    {
        "rows": [
            [("CCP", CONN_YELLOW, BLACK), ("Communists", CONN_YELLOW, BLACK),
             ("Gongchandang", CONN_YELLOW, BLACK), ("The Chinese\nCommunist Party", CONN_YELLOW, BLACK)],
            [("Nationalists", CONN_GREEN, BLACK), ("GMD", CONN_GREEN, BLACK),
             ("Kuomintang", CONN_GREEN, BLACK), ("Guomindang", CONN_GREEN, BLACK)],
            [("Land-owners", CONN_BLUE, BLACK), ("Capitalists", CONN_BLUE, BLACK),
             ("Foreigners", CONN_BLUE, BLACK), ("Gangsters", CONN_BLUE, BLACK)],
            [("Comintern", CONN_UNSOLVED, TEXT_PRIMARY), ("Peasants", CONN_UNSOLVED, TEXT_PRIMARY),
             ("Workers", CONN_UNSOLVED, TEXT_PRIMARY), ("Soviet Union", CONN_UNSOLVED, TEXT_PRIMARY)]
        ]
    },
    # Slide 7: All solved
    {
        "rows": [
            [("CCP", CONN_YELLOW, BLACK), ("Communists", CONN_YELLOW, BLACK),
             ("Gongchandang", CONN_YELLOW, BLACK), ("The Chinese\nCommunist Party", CONN_YELLOW, BLACK)],
            [("Nationalists", CONN_GREEN, BLACK), ("GMD", CONN_GREEN, BLACK),
             ("Kuomintang", CONN_GREEN, BLACK), ("Guomindang", CONN_GREEN, BLACK)],
            [("Land-owners", CONN_BLUE, BLACK), ("Capitalists", CONN_BLUE, BLACK),
             ("Foreigners", CONN_BLUE, BLACK), ("Gangsters", CONN_BLUE, BLACK)],
            [("Peasants", CONN_PURPLE, BLACK), ("Workers", CONN_PURPLE, BLACK),
             ("Comintern", CONN_PURPLE, BLACK), ("Soviet Union", CONN_PURPLE, BLACK)]
        ]
    },
]

for state_idx, state in enumerate(connections_states):
    slide_num = state_idx + 3
    print(f"Building Slide {slide_num}: Connections ({state_idx + 1}/5)...")
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_SLIDE)
    add_accent_bar(slide, Inches(0), SLIDE_W, ACCENT_RED)

    # Title on left
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(2.5), Inches(1.5),
                 "Warm-up:\nConnections!", font_size=32, font_color=TEXT_PRIMARY,
                 bold=True, font_name='Georgia')

    if state_idx == 0:
        add_text_box(slide, Inches(0.5), Inches(3.8), Inches(2.5), Inches(0.5),
                     "Group the 16 terms into 4 rows of 4", font_size=12,
                     font_color=TEXT_SECONDARY, italic=True)

    # Grid
    for row_idx, row in enumerate(state["rows"]):
        for col_idx, (text, bg, fg) in enumerate(row):
            add_conn_cell(slide, col_idx, row_idx, text, bg, fg)


# ═══════════════ SLIDE 8: FOCUS QUESTION / PRC FOUNDING ═══════════════
print("Building Slide 8: Focus Question / PRC Founding...")
slide8 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide8, BG_SLIDE)
add_accent_bar(slide8, Inches(0), SLIDE_W, ACCENT_RED)

add_text_box(slide8, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
             "Founding of the People's Republic of China",
             font_size=28, font_color=TEXT_PRIMARY, bold=True, font_name='Georgia')
add_gold_underline(slide8, Inches(0.7), Inches(1.1))

# Focus Question box
fq_bg = add_card(slide8, Inches(0.7), Inches(1.5), Inches(5.5), Inches(1.4),
                 RGBColor(0x30, 0x22, 0x22), ACCENT_RED)
# Red left border
add_card(slide8, Inches(0.7), Inches(1.5), Pt(5), Inches(1.4), ACCENT_RED)

add_text_box(slide8, Inches(1.0), Inches(1.6), Inches(5.0), Inches(0.3),
             "FOCUS QUESTION", font_size=12, font_color=ACCENT_RED_LT, bold=True)
add_text_box(slide8, Inches(1.0), Inches(1.95), Inches(5.0), Inches(0.8),
             "What factors contributed to the success of the CCP and the downfall of the GMD?",
             font_size=15, font_color=TEXT_PRIMARY, italic=True, font_name='Georgia')

# Image
try:
    img_data2 = download_image(IMG_MAO_PROCLAIM)
    slide8.shapes.add_picture(img_data2, Inches(0.7), Inches(3.2), Inches(5.5), Inches(3.5))
    add_text_box(slide8, Inches(0.7), Inches(6.75), Inches(5.5), Inches(0.3),
                 "Mao Zedong proclaiming the PRC, Tiananmen, Oct. 1 1949 — Wikimedia Commons (Public Domain)",
                 font_size=8, font_color=TEXT_SECONDARY, italic=True)
except Exception as e:
    print(f"  Warning: Could not add image: {e}")

# Right side content
right_body = slide8.shapes.add_textbox(Inches(6.7), Inches(1.5), Inches(6.0), Inches(5.5))
tf = right_body.text_frame
tf.word_wrap = True

# First bullet
p = tf.paragraphs[0]
p.text = "October 1st 1949: Founding of People's Republic of China (PRC)"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_PRIMARY
p.font.name = 'Calibri'
p.space_after = Pt(6)

p2 = tf.add_paragraph()
p2.text = '   "Liberation" (解放) – Mao Zedong gave a triumphant victory speech in Beijing while crowds cheered and watched a procession of the PLA and other party members.'
p2.font.size = Pt(13)
p2.font.color.rgb = TEXT_SECONDARY
p2.font.name = 'Calibri'
p2.level = 1
p2.space_after = Pt(12)

p3 = tf.add_paragraph()
p3.text = "Chiang fled with remaining GMD army to Taiwan to found the Republic of China (ROC)"
p3.font.size = Pt(16)
p3.font.color.rgb = TEXT_PRIMARY
p3.font.name = 'Calibri'
p3.space_after = Pt(6)

for sub_text in [
    "   Claimed to be the legitimate Chinese government",
    "   Until 1971, Taiwan was recognized as the only legitimate government of China by the United Nations"
]:
    p4 = tf.add_paragraph()
    p4.text = sub_text
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_SECONDARY
    p4.font.name = 'Calibri'
    p4.level = 1
    p4.space_after = Pt(4)


# ═══════════════ SLIDE 9: IMPACT OF WWII ═══════════════
print("Building Slide 9: Impact of WWII...")
slide9 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide9, BG_SLIDE)
add_accent_bar(slide9, Inches(0), SLIDE_W, ACCENT_RED)

# Left red panel
red_panel = add_card(slide9, Inches(0.5), Inches(1.2), Inches(3.5), Inches(4.5), ACCENT_RED)
add_text_box(slide9, Inches(0.7), Inches(2.5), Inches(3.1), Inches(2.0),
             "The impact of\nWWII on the\nChinese\nCivil War",
             font_size=26, font_color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Georgia')

# Right content
wwii_body = slide9.shapes.add_textbox(Inches(4.5), Inches(1.5), Inches(8.3), Inches(5.0))
tf = wwii_body.text_frame
tf.word_wrap = True
add_bullet_list(tf, [
    'CCP claimed they more effectively defended China from Japan, while the KMT efforts were described as "half-hearted"',
    "KMT received financial/military support from USA and Great Britain, which gave Chiang legitimacy, but also made him look like a sell out to the West",
    "Chiang expected a long land war with US troops to buy him time to solidify his position against the CCP, which didn't happen after the atomic bomb."
], font_size=18, font_color=TEXT_SECONDARY)


# ═══════════════ SLIDE 10: CIVIL WAR RESUMED ═══════════════
print("Building Slide 10: Civil War Resumed...")
slide10 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide10, BG_SLIDE)
add_accent_bar(slide10, Inches(0), SLIDE_W, ACCENT_RED)

add_text_box(slide10, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
             "The Chinese Civil War Resumed",
             font_size=28, font_color=TEXT_PRIMARY, bold=True, font_name='Georgia')
add_gold_underline(slide10, Inches(0.7), Inches(1.1))

# Left content
resumed_body = slide10.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.0))
tf = resumed_body.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Between 1946-47, the Nationalists made gains in cities and the north and pushed the Communists out of Yan'an."
p.font.size = Pt(16)
p.font.color.rgb = TEXT_SECONDARY
p.font.name = 'Calibri'
p.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "In 1946, Mao reorganized communist forces in a single army People's Liberation Army (PLA)"
p2.font.size = Pt(16)
p2.font.color.rgb = TEXT_SECONDARY
p2.font.name = 'Calibri'
p2.space_after = Pt(4)

p3 = tf.add_paragraph()
p3.text = "   Developed successful guerrilla tactics and took control of Manchuria"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_SECONDARY
p3.font.name = 'Calibri'
p3.level = 1
p3.space_after = Pt(10)

p4 = tf.add_paragraph()
p4.text = "By mid-1947, the PLA turned to conventional warfare and pushed into Nationalist strongholds in the center and west."
p4.font.size = Pt(16)
p4.font.color.rgb = TEXT_SECONDARY
p4.font.name = 'Calibri'

# Map image
try:
    img_data3 = download_image(IMG_CCW_MAP)
    slide10.shapes.add_picture(img_data3, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.0))
    add_text_box(slide10, Inches(7.0), Inches(6.4), Inches(5.8), Inches(0.3),
                 "Chinese Civil War map — Wikimedia Commons (U.S. Government, Public Domain)",
                 font_size=8, font_color=TEXT_SECONDARY, italic=True)
except Exception as e:
    print(f"  Warning: Could not add map: {e}")


# ═══════════════ SLIDE 11: DEVELOPING A COMPLEX THESIS ═══════════════
print("Building Slide 11: Complex Thesis...")
slide11 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide11, BG_SLIDE)
add_accent_bar(slide11, Inches(0), SLIDE_W, ACCENT_RED)

add_text_box(slide11, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
             "Developing a Complex Thesis",
             font_size=28, font_color=TEXT_PRIMARY, bold=True, font_name='Georgia')
add_gold_underline(slide11, Inches(0.7), Inches(1.1))

thesis_body = slide11.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
tf = thesis_body.text_frame
tf.word_wrap = True

items = [
    (0, "A thesis statement is a specific roadmap to a historical argument, not a short-answer response", 16, TEXT_PRIMARY),
    (1, "An evaluation or overall answer to the question", 14, TEXT_SECONDARY),
    (1, "A line of reasoning that shows the moves you will make in your essay", 14, TEXT_SECONDARY),
    (2, "Each point in your line of reasoning should be a category of evidence, not a single piece of evidence (think broad buckets that contain pieces of evidence INSIDE)", 13, TEXT_SECONDARY),
    (0, "Thesis statements are created by first analyzing and organizing evidence into buckets", 16, TEXT_PRIMARY),
    (1, "Group your evidence into groups, then name the category and develop a claim about that group", 14, TEXT_SECONDARY),
    (1, "Once you develop your specific groups, you can see if you have complexity", 14, TEXT_SECONDARY),
    (2, "Multiple groups about a PIECES category, positive or negative categories, continuities or changes, similarities or differences, or some other division that fits your question", 13, TEXT_SECONDARY),
]

for i, (level, text, size, color) in enumerate(items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(3)


# ═══════════════ SLIDE 12: ACTIVITY ═══════════════
print("Building Slide 12: Activity...")
slide12 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide12, BG_SLIDE)
add_accent_bar(slide12, Inches(0), SLIDE_W, ACCENT_RED)

# Central card
act_card = add_card(slide12, Inches(1.5), Inches(1.0), Inches(10.3), Inches(5.5),
                    RGBColor(0x28, 0x28, 0x30), RGBColor(0x3A, 0x3A, 0x42))

add_text_box(slide12, Inches(2.0), Inches(1.3), Inches(9.3), Inches(1.2),
             "What factors contributed to the success of the CCP and downfall of the KMT in Mainland China?",
             font_size=22, font_color=TEXT_PRIMARY, bold=True, font_name='Georgia')

# Red underline
add_card(slide12, Inches(2.0), Inches(2.5), Inches(9.3), Pt(3), ACCENT_RED)

act_body = slide12.shapes.add_textbox(Inches(2.0), Inches(2.8), Inches(9.3), Inches(3.5))
tf = act_body.text_frame
tf.word_wrap = True

act_items = [
    (0, "Work with your table to organize the evidence into three or four groups.", 16, TEXT_SECONDARY),
    (0, "Once you organize the evidence, paraphrase the evidence on the handout and create a claim for each group. Be sure your claim includes a name for that group.", 16, TEXT_SECONDARY),
    (1, "The [GMD/CCP] [won/lost] because…", 14, TEXT_SECONDARY),
    (1, "As you write your claims, think about next-level organization — are your claims political (including military), economic, social?", 14, TEXT_SECONDARY),
    (0, "When you're finished, I will give you a specific prompt that you can respond to with a thesis statement.", 16, TEXT_SECONDARY),
]

for i, (level, text, size, color) in enumerate(act_items):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.level = level
    p.space_before = Pt(4)
    p.space_after = Pt(6)


# ═══════════════ SLIDE 13: WORK TIME ═══════════════
print("Building Slide 13: Work Time...")
slide13 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide13, BG_SLIDE)
add_accent_bar(slide13, Inches(0), SLIDE_W, ACCENT_RED)

# Left panel
wt_panel = add_card(slide13, Inches(0.5), Inches(1.2), Inches(3.0), Inches(5.0),
                    RGBColor(0x2C, 0x2C, 0x34), RGBColor(0x3A, 0x3A, 0x42))
add_text_box(slide13, Inches(0.7), Inches(3.0), Inches(2.6), Inches(1.5),
             "Work\nTime", font_size=34, font_color=TEXT_PRIMARY, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Georgia')

# Task cards
wt_items = [
    (WT_ORANGE, "📄", "You have the rest of class to work on your historical fiction project."),
    (WT_GRAY,   "📋", "Use the Rough Draft Templates for planning and to get teacher feedback."),
    (WT_GOLD,   "💡", "Take a look at some sample scrapbooks or storybooks if you need some ideas."),
    (WT_BLUE,   "📅", "Make sure you have your source for Entry 5 at the start of next class for the reassessment."),
    (WT_GREEN,  "🏫", "If you do not submit a draft, I will assign you to the Second Chance room for the following day. Plan ahead!"),
]

card_top = Inches(1.2)
card_height = Inches(0.85)
card_gap = Inches(0.15)

for i, (color, icon, text) in enumerate(wt_items):
    y = card_top + i * (card_height + card_gap)
    card = add_card(slide13, Inches(4.0), y, Inches(8.8), card_height, color)
    
    # Use text_color based on gold card
    tc = BLACK if color == WT_GOLD else WHITE
    
    add_text_box(slide13, Inches(4.3), y + Inches(0.05), Inches(0.5), card_height,
                 icon, font_size=18)
    add_text_box(slide13, Inches(4.9), y + Inches(0.1), Inches(7.6), Inches(0.7),
                 text, font_size=14, font_color=tc, font_name='Calibri')


# ═══════════════ SAVE ═══════════════
print(f"\nSaving to: {OUTPUT_PATH}")
prs.save(OUTPUT_PATH)
print(f"✅ PowerPoint saved: {OUTPUT_PATH}")
